from pptx import Presentation
from pptx.util import Inches, Pt
import os

prs = Presentation()

slides = [
    ("Barber web — Onlayn barbershop platformasi", "Mijozlar uchun tez rezervatsiya, barberslar uchun boshqaruv, admin uchun nazorat.\nPresenter: (Ismingiz)"),
    ("Kirish", "Muammo: An'anaviy rezervatsiya qiyinchiliklari — qo'ng'iroq, navbatda kutish, xatoliklar.\nYechim: Veb va mobil orqali 1 daqiqada rezervatsiya, avtomatik eslatmalar va boshqaruv paneli."),
    ("Maqsad va vazifalar", "Maqsad: Rezervatsiya jarayonini soddalashtirish.\nVazifalar: Ro'yxatga olish, rezervatsiya, bildirishnomalar, chegirmalar, profiling boshqaruvi."),
    ("Maqsadli auditoriya", "Mijozlar: mobil-first.\nBarberslar: jadval va narxlarni boshqarish.\nAdminlar: tizim nazorat va marketing."),
    ("Asosiy funksiyalar", "- Ro'yxatga olish va login (JWT)\n- Rezervatsiya: barber, xizmat, vaqt, narx\n- Barber paneli: buyurtmalar\n- Admin paneli: chegirmalar\n- Telegram/eslatmalar"),
    ("Foydalanuvchi oqimlari", "Mijoz: kirish → barber → vaqt → tasdiq → eslatma.\nBarber: buyurtma qabul/qaytarish.\nAdmin: chegirma va analytics."),
    ("Texnik arxitektura va stack", "Frontend: React + TypeScript + Vite.\nBackend: Python (FastAPI), SQLAlchemy.\nDB: PostgreSQL.\nMobile: React Native.\nArxitektura: Frontend ↔ API ↔ DB; bot integratsiyasi.") ,
    ("Backend va DB", "`backend/app/` — modelllar, servislar, marshrutlar.\nMigratsiyalar: `backend/alembic/versions/`.\nUploads: `uploads/`.\nAutentifikatsiya: JWT.") ,
    ("Frontend va mobil ilova", "Sahifalar: `src/pages/` ichida.\nAPI: `src/lib/api.ts`.\nMobil: `mobile/App.tsx`.") ,
    ("Bot va integratsiyalar", "Telegram bot orqali tasdiq va eslatmalar.\nWebhook yoki polling orqali backend bilan bog'lanadi.\nBot mijoz va barberni xabardor qiladi.") ,
    ("Joylashtirish va demo", "Backend: .env sozlash, DB bog'lash, `run-local.ps1`.\nFrontend: `npm install` → `npm run dev`.\nMobile: Expo bilan ishga tushirish.") ,
    ("Qulayliklari", "- Tez rezervatsiya\n- No-show kamayadi\n- Barberslar samaraliroq\n- Admin uchun marketing vositalari") ,
    ("Kelgusidagi takomillashtirishlar", "- To'lov integratsiyasi (Stripe)\n- Real-time yangilanishlar\n- Analytics va localization") ,
    ("Savollar va kontakt", "Demo kerakmi? Kontakt: (Email yoki Git repo).")
]

for title, body in slides:
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = body
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(18)

out_dir = os.path.join(os.path.dirname(__file__), "out")
os.makedirs(out_dir, exist_ok=True)
out_path = os.path.join(out_dir, "Barber_web_presentation.pptx")
prs.save(out_path)
print("Created:", out_path)
